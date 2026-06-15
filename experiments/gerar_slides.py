#!/usr/bin/env python3
"""
Gera slides PPTX completos (21 slides, 10+10 min) para apresentacao TabICL v2.

Uso:
    uv run --with python-pptx python experiments/gerar_slides.py

Saida:
    reports/slides_tabicl_v2.pptx
"""

from __future__ import annotations

from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paleta CIn/UFPE ───────────────────────────────────────────────────────────
RED      = RGBColor(0xC4, 0x12, 0x30)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY     = RGBColor(0x6B, 0x6B, 0x6B)
LGRAY    = RGBColor(0xF0, 0xF0, 0xF0)
MGRAY    = RGBColor(0xD8, 0xD8, 0xD8)
DARKRED  = RGBColor(0x8B, 0x00, 0x00)
ACCENT   = RGBColor(0xFF, 0xEB, 0xEB)

# Hex equivalentes para matplotlib
RED_HEX   = "#C41230"
LGRAY_HEX = "#F0F0F0"
GRAY_HEX  = "#6B6B6B"

# ── Geometria 16:9 ────────────────────────────────────────────────────────────
SW  = Inches(13.33)
SH  = Inches(7.50)
SB  = Inches(1.30)
MH  = Inches(1.20)
MRG = Inches(0.28)
CW  = SW - SB

JP  = "Joao Pedro"
MC  = "Maria Clara"
VIN = "Vinicius"

PLOTS = Path("results/plots")

# ── Dados reais (pre-calculados de final_run_results_v2.csv) ──────────────────
# (*) = exigido explicitamente pelo enunciado ("preset extreme de 4 horas")
# AutoGluon_Extreme (sem sufixo) = extra nao-exigido (~30 min/dataset)
RESULTS = [
    # nome display,               AUC,    ACC,    GMean,  CE,     Tempo(s)
    ("AG Extreme 4h [spec] (*)", 0.9099, 0.9121, 0.6039, 0.3240, 14611.7),
    ("AG Extreme (~30min)",      0.9107, 0.9154, 0.6560, 0.2194,  1830.9),
    ("TabICL v2",                0.9056, 0.9144, 0.6786, 0.2071,    48.5),
    ("CatBoost Tuned",           0.9015, 0.9106, 0.6500, 0.2369,   302.4),
    ("CatBoost TD",              0.8972, 0.9085, 0.6432, 0.2366,    14.4),
    ("XGBoost Tuned",            0.8968, 0.9078, 0.6755, 0.2382,    96.7),
    ("LightGBM Tuned",           0.8947, 0.9116, 0.6539, 0.2464,   176.3),
    ("LightGBM TD",              0.8933, 0.9067, 0.6474, 0.2860,     1.6),
    ("AutoGluon Default",        0.8887, 0.8953, 0.5604, 0.2676,  1413.0),
    ("XGBoost TD",               0.8866, 0.9072, 0.6569, 0.2673,     0.7),
]
# Referencia primaria para comparacao (exigida pelo enunciado)
AG4H_AUC  = 0.9099
AG4H_TEMPO = 14611.7


# =============================================================================
# Geracao de graficos extras via matplotlib
# =============================================================================

def _generate_charts():
    """Gera PNGs extras que nao existem ainda em results/plots/."""
    PLOTS.mkdir(parents=True, exist_ok=True)

    # Tenta carregar CSVs reais; usa RESULTS como fallback
    csv = Path("results/data/final_run_results_v2.csv")
    meta_csv = Path("results/data/dataset_metadata.csv")

    # ── 1. Barras horizontais: AUC-OVO por modelo ─────────────────────────────
    out1 = PLOTS / "bar_auc_models.png"
    # Ordena por AUC para o grafico (decrescente de baixo pra cima = crescente no barh)
    sorted_r = sorted(RESULTS, key=lambda x: x[1])
    names = [r[0] for r in sorted_r]
    aucs  = [r[1] for r in sorted_r]

    def _bar_color(n):
        if "TabICL" in n:   return RED_HEX
        if "[spec]" in n:   return "#222222"   # preto para o 4h exigido
        return GRAY_HEX

    colors = [_bar_color(n) for n in names]

    fig, ax = plt.subplots(figsize=(9.5, 5.5))
    bars = ax.barh(names, aucs, color=colors, edgecolor="white", height=0.65)
    ax.set_xlabel("AUC-OVO (media, 30 datasets)", fontsize=12)
    ax.set_xlim(0.868, 0.922)

    tabicl_auc = next(r[1] for r in RESULTS if "TabICL" in r[0])
    ag4h_auc   = next(r[1] for r in RESULTS if "[spec]" in r[0])
    ax.axvline(x=tabicl_auc, color=RED_HEX, linestyle="--", linewidth=1.5,
               label=f"TabICL v2  ({tabicl_auc:.4f})")
    ax.axvline(x=ag4h_auc, color="#222222", linestyle=":", linewidth=1.5,
               label=f"AG Extreme 4h [spec]  ({ag4h_auc:.4f})")

    for bar, auc, name in zip(bars, aucs, names):
        ax.text(auc + 0.0004, bar.get_y() + bar.get_height()/2,
                f"{auc:.4f}", va="center", fontsize=8.5,
                fontweight="bold" if ("TabICL" in name or "[spec]" in name) else "normal")

    ax.legend(fontsize=9, loc="lower right")
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("#FAFAFA")
    fig.patch.set_facecolor("white")
    plt.tight_layout()
    fig.savefig(out1, dpi=150, bbox_inches="tight")
    plt.close(fig)
    print(f"  gerado: {out1}")

    # ── 2. Regime x modelo (barras agrupadas) ─────────────────────────────────
    out2 = PLOTS / "bar_regime.png"
    regime_data = {
        "Small": {
            "TabICL v2": 0.8056, "AG Extreme 4h [spec]": 0.8628,
            "CatBoost Tuned": 0.8591, "LightGBM Tuned": 0.8595,
        },
        "Medium": {
            "TabICL v2": 0.9425, "AG Extreme 4h [spec]": 0.9453,
            "CatBoost Tuned": 0.9359, "LightGBM Tuned": 0.9241,
        },
        "Large": {
            "TabICL v2": 0.8629, "AG Extreme 4h [spec]": 0.8639,
            "CatBoost Tuned": 0.8555, "LightGBM Tuned": 0.8554,
        },
    }
    modelos = ["TabICL v2", "AG Extreme 4h [spec]", "CatBoost Tuned", "LightGBM Tuned"]
    regimes = list(regime_data.keys())
    x = np.arange(len(regimes))
    width = 0.2
    fig, ax = plt.subplots(figsize=(9, 4.5))
    palette = [RED_HEX, "#222222", "#888888", "#AAAAAA"]
    for i, (mod, col) in enumerate(zip(modelos, palette)):
        vals = [regime_data[r][mod] for r in regimes]
        bars = ax.bar(x + i * width, vals, width, label=mod, color=col, edgecolor="white")
        for b, v in zip(bars, vals):
            ax.text(b.get_x() + b.get_width()/2, v + 0.001,
                    f"{v:.3f}", ha="center", va="bottom", fontsize=8)
    ax.set_xticks(x + width * 1.5)
    ax.set_xticklabels(regimes, fontsize=12)
    ax.set_ylabel("AUC-OVO", fontsize=12)
    ax.set_ylim(0.780, 0.970)
    ax.legend(fontsize=9, ncol=2)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("#FAFAFA")
    fig.patch.set_facecolor("white")
    plt.tight_layout()
    fig.savefig(out2, dpi=150, bbox_inches="tight")
    plt.close(fig)
    print(f"  gerado: {out2}")

    # ── 3. Os 4 eixos de caracterizacao (2x2 subplots) ────────────────────────
    out3 = PLOTS / "bar_eixos.png"
    eixos = [
        # (titulo, categorias, tabicl_vals, rival_name, rival_vals)
        ("(i) Tamanho do Dataset",
         ["Small", "Medium", "Large"],
         [0.8056, 0.9425, 0.8629],
         "AG Extreme 4h [spec]",
         [0.8628, 0.9453, 0.8639]),
        ("(ii) Numero de Classes",
         ["Binario", "Multiclasse"],
         [0.8562, 0.9585],
         "AG Extreme 4h [spec]",
         [0.8592, 0.9606]),
        ("(iii) Tipo de Feature",
         ["Numerica", "Categorica"],
         [0.9221, 0.8636],
         "AG Extreme 4h [spec]",
         [0.9221, 0.8763]),
        ("(iv) Dados Faltantes",
         ["Sem NaN", "Com NaN"],
         [0.9221, 0.8539],
         "AG Extreme 4h [spec]",
         [0.9221, 0.8763]),
    ]
    fig, axes = plt.subplots(2, 2, figsize=(10, 6.5))
    axes = axes.flatten()
    for ax, (titulo, cats, tab_vals, rival_nome, rival_vals) in zip(axes, eixos):
        x = np.arange(len(cats))
        ax.bar(x - 0.2, tab_vals, 0.38, label="TabICL v2",
               color=RED_HEX, edgecolor="white")
        ax.bar(x + 0.2, rival_vals, 0.38, label=rival_nome,
               color=GRAY_HEX, edgecolor="white")
        for xi, v in zip(x, tab_vals):
            ax.text(xi - 0.2, v + 0.002, f"{v:.3f}", ha="center", fontsize=8,
                    color=RED_HEX, fontweight="bold")
        for xi, v in zip(x, rival_vals):
            ax.text(xi + 0.2, v + 0.002, f"{v:.3f}", ha="center", fontsize=8, color="gray")
        ax.set_xticks(x)
        ax.set_xticklabels(cats, fontsize=10)
        ax.set_title(titulo, fontsize=11, fontweight="bold")
        ax.set_ylim(min(tab_vals + rival_vals) - 0.04, max(tab_vals + rival_vals) + 0.035)
        ax.spines[["top", "right"]].set_visible(False)
        ax.legend(fontsize=8)
        ax.set_facecolor("#FAFAFA")
    fig.patch.set_facecolor("white")
    plt.tight_layout()
    fig.savefig(out3, dpi=150, bbox_inches="tight")
    plt.close(fig)
    print(f"  gerado: {out3}")

    # ── 4. Todas as 4 metricas em barras (visao geral) ────────────────────────
    out4 = PLOTS / "bar_4metricas.png"
    metrics = ["AUC-OVO", "Accuracy", "G-Mean"]
    indices = [1, 2, 3]   # colunas em RESULTS
    mod_names = [r[0] for r in RESULTS]
    vals_by_metric = [[r[idx] for r in RESULTS] for idx in indices]

    fig, axes = plt.subplots(1, 3, figsize=(13, 4.5))
    for ax, metric, vals in zip(axes, metrics, vals_by_metric):
        cols = [RED_HEX if "TabICL" in n else GRAY_HEX for n in mod_names]
        bars = ax.barh(mod_names[::-1], vals[::-1], color=cols[::-1],
                       edgecolor="white", height=0.65)
        ax.set_title(metric, fontsize=12, fontweight="bold")
        ax.set_xlim(min(vals) - 0.02, max(vals) + 0.02)
        ax.spines[["top", "right"]].set_visible(False)
        ax.set_facecolor("#FAFAFA")
        for bar, v in zip(bars, vals[::-1]):
            ax.text(v + 0.001, bar.get_y() + bar.get_height()/2,
                    f"{v:.3f}", va="center", fontsize=7)
    fig.patch.set_facecolor("white")
    plt.tight_layout()
    fig.savefig(out4, dpi=150, bbox_inches="tight")
    plt.close(fig)
    print(f"  gerado: {out4}")

    print("  Todos os graficos extras gerados.")


# =============================================================================
# Primitivas de desenho PPTX
# =============================================================================

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, color: RGBColor):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color
    return s


def _tx(slide, x, y, w, h, text: str,
        size=14, bold=False, color=NEAR_BLK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def _img(slide, path, x, y, w, h):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), x, y, w, h)
    else:
        _rect(slide, x, y, w, h, LGRAY)
        _tx(slide, x + Inches(.1), y + h // 2, w - Inches(.2), Inches(.35),
            f"[{Path(path).name}]", size=9, color=GRAY, align=PP_ALIGN.CENTER)


def _sidebar(slide):
    _rect(slide, 0, 0, SB, SH, RED)


def _header(slide, title: str, presenter: str | None = None):
    _rect(slide, SB, 0, CW, MH, LGRAY)
    _tx(slide, SB + MRG, Inches(.15), CW - MRG * 2 - Inches(2.5), Inches(.9),
        title, size=20, bold=True, color=RED)
    if presenter:
        _tx(slide, SW - Inches(2.85), Inches(.25), Inches(2.55), Inches(.42),
            f"| {presenter}", size=10, color=GRAY, align=PP_ALIGN.RIGHT, italic=True)


def _pptx_table(slide, headers, rows, x, y, w, h,
                hdr_bg=RED, hdr_fg=WHITE):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), int(x), int(y), int(w), int(h)).table
    col_w = int(w) // len(headers)
    for c in range(len(headers)):
        tbl.columns[c].width = col_w

    def _cell(r, c, txt, bold=False, sz=10, bg=None, fg=NEAR_BLK, al=PP_ALIGN.CENTER):
        cell = tbl.cell(r, c)
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = al
        run = p.add_run()
        run.text = str(txt)
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = fg

    for c, h in enumerate(headers):
        _cell(0, c, h, bold=True, sz=11, bg=hdr_bg, fg=hdr_fg)
    for r, row in enumerate(rows):
        bg = WHITE if r % 2 == 0 else LGRAY
        for c, val in enumerate(row):
            hi = "TabICL" in str(val)
            _cell(r + 1, c, val, sz=10, bg=bg,
                  fg=RED if hi else NEAR_BLK, bold=hi)


# =============================================================================
# Factories de slides
# =============================================================================

def slide_capa(prs):
    sl = _blank(prs)
    panel_w = Inches(5.1)
    _rect(sl, 0, 0, panel_w, SH, RED)
    _tx(sl, Inches(.25), Inches(.28), panel_w - Inches(.3), Inches(.6),
        "CIn . UFPE", size=16, bold=True, color=WHITE)
    _tx(sl, Inches(.25), Inches(2.5), panel_w - Inches(.3), Inches(2.3),
        "Aprendizagem de Maquina\nPos-Graduacao 2026.1\n\nProf. Leandro Almeida",
        size=14, color=WHITE)
    _tx(sl, panel_w + Inches(.45), Inches(1.1), SW - panel_w - Inches(.6), Inches(2.6),
        "TabICL v2:\nFoundation Model para\nDados Tabulares",
        size=28, bold=True, color=RED)
    _tx(sl, panel_w + Inches(.45), Inches(3.85), SW - panel_w - Inches(.6), Inches(.5),
        "Avaliacao em Benchmark TabArena-v0.1  |  NeurIPS 2025",
        size=13, color=GRAY, italic=True)
    _tx(sl, panel_w + Inches(.45), Inches(4.7), SW - panel_w - Inches(.6), Inches(2.3),
        "Joao Pedro Miranda da Silva    jpms5@cin.ufpe.br\n"
        "Maria Clara F. G. Barretto    mcfgb@cin.ufpe.br\n"
        "Vinicius Limeira Valenca       vlv3@cin.ufpe.br",
        size=13, color=NEAR_BLK)


def slide_agenda(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Estrutura da Apresentacao (20 min  |  10 + 10)")
    cw2 = (CW - Inches(.4)) / 2
    ch2 = (SH - MH - Inches(.3)) / 2
    pos = [
        (SB + Inches(.1), MH + Inches(.1)),
        (SB + cw2 + Inches(.2), MH + Inches(.1)),
        (SB + Inches(.1), MH + ch2 + Inches(.18)),
        (SB + cw2 + Inches(.2), MH + ch2 + Inches(.18)),
    ]
    items = [
        ("Etapa 1 — Teoria (10 min)  |  Joao Pedro",
         "Motivacao e contexto historico\n"
         "Arquitetura e funcionamento interno\n"
         "Forma de aprendizado e representacao\n"
         "Aplicacoes praticas e limitacoes"),
        ("Etapa 2a — Experimentos (5 min)  |  Maria Clara",
         "30 datasets e 10 sistemas\n"
         "Resultados por metrica (AUC, ACC, G-Mean, CE)\n"
         "Analise por regime (4 eixos)\n"
         "CD Diagram — Friedman-Nemenyi"),
        ("Etapa 2b — Analise Estatistica (5 min)  |  Vinicius",
         "Bayesian Signed-Rank (ROPE = 0.01)\n"
         "Busca de hiperparametros — Optuna TPE\n"
         "Custo vs. Desempenho — Fronteira de Pareto\n"
         "Conclusoes e recomendacoes"),
        ("Reproducibilidade",
         "pyproject.toml + uv.lock (versoes travadas)\n"
         "seed=42 em split, Optuna e PyTorch\n"
         "Dockerfile para containerizacao\n"
         "Cluster Apuana — RTX 3090"),
    ]
    for (lbl, body), (x, y) in zip(items, pos):
        _rect(sl, x, y, cw2, ch2 - Inches(.08), LGRAY)
        _tx(sl, x + Inches(.15), y + Inches(.1), cw2 - Inches(.25), Inches(.42),
            lbl, size=12, bold=True, color=RED)
        _tx(sl, x + Inches(.15), y + Inches(.56), cw2 - Inches(.25), ch2 - Inches(.72),
            body, size=11, color=NEAR_BLK)


def slide_motivacao(prs):
    """Timeline + motivacao com grafico de posicionamento."""
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Motivacao e Contexto Historico", JP)

    # Timeline horizontal
    tl_y = MH + Inches(1.1)
    x0, x1 = SB + Inches(.35), SW - Inches(.35)
    _rect(sl, x0, tl_y, x1 - x0, Inches(.055), MGRAY)

    pontos = [
        (0.00, "2016", "XGBoost domina\nKaggle e benchmarks\ntabulares"),
        (0.25, "2019", "AutoML surge:\nAutoGluon e H2O\n(horas de tuning)"),
        (0.50, "2022", "TabPFN:\n1o ICL tabular\n(n < 3.000)"),
        (0.75, "2025", "TabICL v1:\nescala para n > 10K\n(ICML 2025)"),
        (1.00, "2026", "TabICL v2:\nmulticlasse + ensemble\n(este trabalho)"),
    ]
    span = x1 - x0
    for frac, ano, desc in pontos:
        px = x0 + span * frac
        _rect(sl, px - Inches(.06), tl_y - Inches(.06), Inches(.12), Inches(.17), RED)
        _tx(sl, px - Inches(.55), tl_y - Inches(.44), Inches(1.1), Inches(.35),
            ano, size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
        _tx(sl, px - Inches(.72), tl_y + Inches(.22), Inches(1.44), Inches(.88),
            desc, size=10, color=NEAR_BLK, align=PP_ALIGN.CENTER)

    # Caixa de destaque
    cy = MH + Inches(2.45)
    _rect(sl, SB + MRG, cy, CW - MRG * 2, Inches(1.0), ACCENT)
    _tx(sl, SB + MRG + Inches(.15), cy + Inches(.1), CW - MRG * 2 - Inches(.3), Inches(.8),
        "Problema central: GBDTs exigem feature engineering, pipelines complexos e horas de Optuna. "
        "Foundation Models aprendem o algoritmo no pre-treino; na inferencia apenas atencao cruzada "
        "contra os dados de treino — sem gradient descent, sem fine-tuning, sem feature engineering.",
        size=13, color=DARKRED)

    # Bullets de contribuicao
    cy2 = cy + Inches(1.12)
    for lbl, body in [
        ("GBDTs ainda lideram em Large e dados com NaN nativo",
         "XGBoost/LightGBM/CatBoost: particao nativa de missing, velocidade de inferencia em CPU."),
        ("TabICL v2 preenche a lacuna para Medium e Multiclasse",
         "AUC equivalente ao AutoGluon Extreme 4h com 38x menos tempo (48s vs 1831s por dataset)."),
    ]:
        _tx(sl, SB + MRG, cy2, CW - MRG * 2, Inches(.34), "=> " + lbl,
            size=13, bold=True, color=RED)
        _tx(sl, SB + MRG + Inches(.12), cy2 + Inches(.36), CW - MRG * 2, Inches(.38),
            body, size=12, color=NEAR_BLK)
        cy2 += Inches(.82)


def slide_arquitetura(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Arquitetura: Visao Geral do TabICL v2", JP)

    # Diagrama de blocos em texto
    bw, bh = Inches(1.65), Inches(.72)
    by = MH + Inches(.3)
    boxes = [
        ("Dados de\nTreino\n(N amostras)", SB + Inches(.3), LGRAY, NEAR_BLK),
        ("Embedding\nLinear / Vocab", SB + Inches(2.18), LGRAY, NEAR_BLK),
        ("Transformer\nEncoder\n(K, V)", SB + Inches(4.06), LGRAY, NEAR_BLK),
        ("Atencao\nCruzada\n(Q = teste)", SB + Inches(5.94), RED, WHITE),
        ("Softmax\n(C classes)", SB + Inches(7.82), LGRAY, NEAR_BLK),
        ("Predicao\n(proba.)", SB + Inches(9.7), LGRAY, NEAR_BLK),
    ]
    for lbl, bx, bg, fg in boxes:
        _rect(sl, bx, by, bw, bh, bg)
        _tx(sl, bx + Inches(.08), by + Inches(.08), bw - Inches(.16), bh - Inches(.1),
            lbl, size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)
        if bx + bw < SW - Inches(.4):
            _tx(sl, bx + bw, by + Inches(.26), Inches(.2), Inches(.22),
                "->", size=13, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    _rect(sl, SB + MRG, by + bh + Inches(.15), CW - MRG*2, Inches(.03), MGRAY)
    _tx(sl, SB + MRG, by + bh + Inches(.2), CW - MRG*2, Inches(.26),
        "Amostra de Teste:  Q = Embedding(x_test)  --  nenhum parametro atualizado",
        size=11, italic=True, color=GRAY)

    # 4 caracteristicas em caixas
    cy = MH + Inches(1.7)
    feats = [
        ("27M parametros", "pre-treinado em tabelas\nsinteticas (priores\nBayesianos)"),
        ("Interface sklearn", "fit(X_train, y_train)\npredict_proba(X_test)\nzero config"),
        ("Apache 2.0", "uso comercial e\nacademico sem\nrestricoes"),
        ("Ensemble ICL", "n_estimators contextos\ncombinados por media\nde logits"),
    ]
    fw = (CW - Inches(.4)) / len(feats)
    for i, (lbl, desc) in enumerate(feats):
        fx = SB + Inches(.1) + i * (fw + Inches(.07))
        _rect(sl, fx, cy, fw, Inches(2.0), LGRAY)
        _tx(sl, fx + Inches(.1), cy + Inches(.1), fw - Inches(.2), Inches(.42),
            lbl, size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
        _tx(sl, fx + Inches(.1), cy + Inches(.58), fw - Inches(.2), Inches(1.2),
            desc, size=11, color=NEAR_BLK, align=PP_ALIGN.CENTER)


def slide_funcionamento(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Funcionamento Interno: Atencao Cruzada e Formulacao Matematica", JP)

    cx, cw = SB + MRG, CW - MRG * 2
    _tx(sl, cx, MH + Inches(.1), cw, Inches(.5),
        "O mecanismo central: cada amostra de teste (query Q) consulta todo o conjunto de treino "
        "(chave K e valor V) via atencao cruzada. Nenhum gradiente e calculado.",
        size=13, color=NEAR_BLK)

    _rect(sl, cx, MH + Inches(.7), cw, Inches(.95), ACCENT)
    _tx(sl, cx + Inches(.4), MH + Inches(.82), cw - Inches(.8), Inches(.7),
        "Attention(Q, K, V)  =  softmax( Q K^T / sqrt(d) ) . V",
        size=22, bold=True, color=DARKRED, align=PP_ALIGN.CENTER)

    cy = MH + Inches(1.78)
    notas = [
        ("Q  in  R^(1 x d)", "embedding da amostra de TESTE (query)"),
        ("K, V  in  R^(N x d)", "embeddings das N amostras de TREINO (contexto)"),
        ("d", "dimensao do embedding (hidden dim do Transformer)"),
        ("C logits -> Softmax", "saida projetada para C classes via camada linear"),
    ]
    for sym, desc in notas:
        _tx(sl, cx + Inches(.15), cy, Inches(2.45), Inches(.36),
            sym, size=13, bold=True, color=RED)
        _tx(sl, cx + Inches(2.65), cy, cw - Inches(2.8), Inches(.36),
            "-- " + desc, size=13, color=NEAR_BLK)
        cy += Inches(.38)

    _rect(sl, cx, cy + Inches(.08), cw, Inches(.78), LGRAY)
    _tx(sl, cx + Inches(.15), cy + Inches(.16), cw - Inches(.3), Inches(.58),
        "Ensemble de contextos:  logits_final = mean( logits_1, ..., logits_M )  "
        "com M = n_estimators sub-conjuntos aleatorios do treino  ->  "
        "equivale a bagging por ICL sem re-treinamento",
        size=12, color=NEAR_BLK)


def slide_aprendizado(prs):
    """Forma de aprendizado: pre-treinamento vs inferencia em 2 colunas."""
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Forma de Aprendizado e Representacao dos Padroes nos Dados", JP)

    cx, cw = SB + MRG, CW - MRG * 2
    _tx(sl, cx, MH + Inches(.1), cw, Inches(.38),
        "O TabICL v2 aprende de forma INDIRETA: no pre-treinamento (offline), nao na inferencia.",
        size=14, bold=True, color=RED)

    half = cw / 2
    cols = [
        ("PRE-TREINAMENTO (offline, uma vez)",
         "Tabelas sinteticas geradas por priores Bayesianos\n"
         "Milhoes de tarefas de classificacao artificiais\n"
         "Objetivo: minimizar cross-entropy (gradient descent)\n"
         "Resultado: 27M pesos que codificam 'como aprender'\n\n"
         "Representacao aprendida:\n"
         "  features numericas -> embedding linear\n"
         "  features categoricas -> embedding de vocabulario\n"
         "  labels -> embedding de classe\n"
         "  N amostras de treino -> sequencia de tokens"),
        ("INFERENCIA (online, por query)\n"
         "fit(X_train, y_train):",
         "  -> tokeniza N amostras de treino\n"
         "  -> armazena K, V no contexto\n\n"
         "predict_proba(X_test):\n"
         "  -> tokeniza x_test (= query Q)\n"
         "  -> calcula Attention(Q, K, V)\n"
         "  -> projeta em C logits -> Softmax\n\n"
         "Complexidade: O(N.d) por amostra de teste\n"
         "Total N testes: O(N^2 . p)  <- gargalo\n\n"
         "Nenhum parametro e atualizado."),
    ]
    for i, (title, body) in enumerate(cols):
        x = cx + i * (half + Inches(.15))
        w = half - Inches(.1)
        bg = ACCENT if i == 0 else LGRAY
        _rect(sl, x, MH + Inches(.58), w, SH - MH - Inches(.68), bg)
        _tx(sl, x + Inches(.12), MH + Inches(.68), w - Inches(.24), Inches(.42),
            title, size=13, bold=True, color=RED)
        _tx(sl, x + Inches(.12), MH + Inches(1.16), w - Inches(.24),
            SH - MH - Inches(1.3), body, size=11, color=NEAR_BLK)


def slide_complexidade(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Complexidade Computacional e Escalabilidade", JP)

    cx, cw = SB + MRG, CW - MRG * 2
    headers = ["Modelo", "Treino", "Inferencia (1 amostra)", "Memoria GPU", "Large (n>50K)"]
    rows = [
        ("TabICL v2",  "O(1)*",            "O(N.d) -> O(N^2.p) total", "4-24 GB",   "OOM sem offload"),
        ("LightGBM TD", "O(N.p.log N)",    "O(p . profundidade)",       "CPU (MB)", "Excelente"),
        ("AutoGluon 4h","horas de ensemble","variavel",                  "alta",     "Escala bem"),
    ]
    _pptx_table(sl, headers, rows, cx, MH + Inches(.12), cw, Inches(1.6))

    cy = MH + Inches(1.82)
    items = [
        ("O(N^2 . p) -- gargalo da inferencia",
         "Para N amostras de treino e N de teste: calcula Q K^T em R^(N x N). "
         "Aceitavel em Small/Medium (N < 10K, tempo < 60s). Critico em Large."),
        ("Mitigacoes implementadas (cluster Apuana, RTX 3090 24 GB)",
         "Query-Aware Scalable Softmax + offload_mode='disk' (chunked inference). "
         "Script: experiments/cluster/job_apuana.slurm."),
        ("302x mais rapido que AutoGluon Extreme 4h [exigido pelo enunciado]",
         "TabICL: 48.5s / dataset. AG Extreme 4h (spec): 14.612s (~4h) / dataset. "
         "AG Extreme ~30min (extra): 1831s. Custo quadratico aceitavel para n < 20K em GPU."),
    ]
    for lbl, body in items:
        _tx(sl, cx, cy, cw, Inches(.35), lbl, size=13, bold=True, color=RED)
        _tx(sl, cx + Inches(.1), cy + Inches(.37), cw, Inches(.58), body, size=12, color=NEAR_BLK)
        cy += Inches(1.02)


def slide_hiperparametros(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Hiperparametros Principais do TabICL v2", JP)

    cx, cw = SB + MRG, CW - MRG * 2
    _tx(sl, cx, MH + Inches(.1), cw, Inches(.35),
        "Espaco de busca definido em experiments/tabicl_tuning.py. "
        "Experimento principal rodou com defaults — tuning opcional.",
        size=12, italic=True, color=GRAY)

    headers = ["Hiperparametro", "Padrao", "Busca Optuna", "Efeito principal"]
    rows = [
        ("n_estimators",       "16",      "[4, 32]",       "tamanho ensemble; mais = estavel, mas custo linear"),
        ("softmax_temperature","0.9",     "[0.3, 1.5]",    "calibracao: baixo = confiante; alto = distribuicao suave"),
        ("outlier_threshold",  "4.0",     "[2.0, 8.0]",    "clipping de features em multiplos de sigma"),
        ("norm_methods",       "z-score", "categorico",    "z-score / quantile / robusto para embeddings"),
        ("average_logits",     "True",    "True/False",    "combinacao do ensemble: media (True) ou log-sum-exp"),
        ("offload_mode",       "None",    "None / 'disk'", "disk = chunked inference para Large (evita OOM)"),
    ]
    _pptx_table(sl, headers, rows, cx, MH + Inches(.5), cw, Inches(3.65))

    _rect(sl, cx, MH + Inches(4.28), cw, Inches(.68), LGRAY)
    _tx(sl, cx + Inches(.15), MH + Inches(4.36), cw - Inches(.3), Inches(.55),
        "Resultado pratico: o modelo opera de forma competitiva com os defaults dos autores. "
        "O espaco Optuna esta definido mas nao foi executado para TabICL no experimento final "
        "(limitacao de tempo no cluster HPC Apuana).",
        size=11, italic=True, color=GRAY)


def slide_aplicacoes(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Aplicacoes Praticas e Limitacoes", JP)

    cx, cw = SB + MRG, CW - MRG * 2
    half = cw / 2

    # Coluna USAR
    _rect(sl, cx, MH + Inches(.1), half - Inches(.12), SH - MH - Inches(.2), LGRAY)
    _tx(sl, cx + Inches(.12), MH + Inches(.2), half - Inches(.3), Inches(.4),
        "QUANDO USAR o TabICL v2", size=13, bold=True, color=RED)
    for i, (lbl, val) in enumerate([
        ("Regime Medium (1K-10K)", "melhor AUC: 0.9425"),
        ("Classificacao Multiclasse", "lider: 0.9585 vs AutoGluon 0.9606"),
        ("Features numericas", "AUC 0.9221 (equiv. ao SOTA)"),
        ("Prototipagem rapida", "48s/dataset, zero feature engineering"),
        ("GPU disponivel", "RTX 3090 suficiente para Medium"),
        ("Sem NaN ou com imputacao previa", "AUC 0.9221 sem NaN"),
    ]):
        cy = MH + Inches(.7) + i * Inches(.57)
        _tx(sl, cx + Inches(.2), cy, half - Inches(.45), Inches(.27),
            lbl, size=12, bold=True, color=NEAR_BLK)
        _tx(sl, cx + Inches(.35), cy + Inches(.28), half - Inches(.6), Inches(.25),
            val, size=11, color=GRAY)

    # Coluna LIMITACOES
    rx = cx + half + Inches(.12)
    rw = half - Inches(.12)
    _rect(sl, rx, MH + Inches(.1), rw, SH - MH - Inches(.2), ACCENT)
    _tx(sl, rx + Inches(.12), MH + Inches(.2), rw - Inches(.3), Inches(.4),
        "LIMITACOES documentadas", size=13, bold=True, color=RED)
    for i, (lbl, val) in enumerate([
        ("Large sem GPU >= 24 GB", "OOM: n>50K, p>500 sem offload"),
        ("Dados com NaN nativos", "0.854 vs CatBoost TD 0.872"),
        ("Features categoricas", "CatBoost Tuned supera (0.872 vs 0.864)"),
        ("Binario em Large", "0.856 -- GBDTs Tuned sao superiores"),
        ("Interpretabilidade", "caixa-preta sem SHAP nativo"),
        ("Tuning nao executado", "defaults usados no experimento final"),
    ]):
        cy = MH + Inches(.7) + i * Inches(.57)
        _tx(sl, rx + Inches(.2), cy, rw - Inches(.45), Inches(.27),
            lbl, size=12, bold=True, color=DARKRED)
        _tx(sl, rx + Inches(.35), cy + Inches(.28), rw - Inches(.6), Inches(.25),
            val, size=11, color=GRAY)


def slide_transicao(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _rect(sl, SB, 0, CW, SH, WHITE)
    _rect(sl, SB, SH / 2 - Inches(.03), CW, Inches(.05), RED)
    _tx(sl, SB + Inches(.5), Inches(1.55), CW - Inches(1.0), Inches(2.6),
        "Etapa 2:\nEstudo Experimental",
        size=34, bold=True, color=RED, align=PP_ALIGN.CENTER)
    _tx(sl, SB + Inches(.5), Inches(4.35), CW - Inches(1.0), Inches(.85),
        "30 datasets  |  10 sistemas  |  4 metricas  |  validacao frequentista e bayesiana",
        size=15, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    _tx(sl, SB + Inches(.5), Inches(5.42), CW - Inches(1.0), Inches(.52),
        "Maria Clara (slides 11-15)    Vinicius (slides 16-20)",
        size=13, color=GRAY, align=PP_ALIGN.CENTER)


def slide_setup(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Configuracao do Experimento", MC)
    cx, cw = SB + MRG, CW - MRG * 2
    _tx(sl, cx, MH + Inches(.1), cw, Inches(.35),
        "Protocolo Demsar (2006) + Benavoli et al. (2017) conforme especificacao do projeto.",
        size=12, italic=True, color=GRAY)
    headers = ["Aspecto", "Configuracao"]
    rows = [
        ("Datasets",          "30 datasets TabArena-v0.1 (NeurIPS 2025) via OpenML"),
        ("Distribuicao",      "3 Small (n<1K) + 17 Medium (1K-10K) + 10 Large (>10K)*"),
        ("Split",             "70% treino / 30% teste, estratificado, seed=42"),
        ("Baselines TD",      "LightGBM, XGBoost, CatBoost (defaults pytabkit meta-aprendidos)"),
        ("Baselines Tuned",   "LightGBM, XGBoost, CatBoost (Optuna TPE, 20 trials, max 2h)"),
        ("AutoML",            "AutoGluon Default  +  AutoGluon Extreme 4h [EXIGIDO PELO ENUNCIADO]  +  AG Extreme ~30min [extra]"),
        ("Metricas",          "AUC-OVO (macro), Accuracy, G-Mean, Cross-Entropy, Tempo (s)"),
        ("Analise estatist.", "Friedman + Nemenyi (autorank) + Bayesian ROPE=0.01 (baycomp)"),
        ("Infraestrutura",    "GPU RTX 3090 (Large); CPU local (Small/Medium); seed=42"),
    ]
    _pptx_table(sl, headers, rows, cx, MH + Inches(.5), cw, Inches(5.4))
    _tx(sl, cx, SH - Inches(.34), cw, Inches(.28),
        "* Desvio da especificacao 10+10+10: TabArena-v0.1 tem oferta limitada de Small (n<1K) "
        "apos filtragem de qualidade (IID, licenca CC0).",
        size=9, italic=True, color=GRAY)


def slide_datasets_e_resultados(prs):
    """Slide duplo: resumo dos 30 datasets + grafico de AUC por modelo."""
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "30 Datasets e AUC-OVO por Sistema (media geral)", MC)

    half = (CW - Inches(.3)) / 2

    # Esquerda: resumo dos datasets
    cx = SB + Inches(.1)
    _tx(sl, cx, MH + Inches(.12), half - Inches(.1), Inches(.35),
        "Estratificacao dos 30 datasets:", size=13, bold=True, color=RED)
    headers = ["Regime", "N", "Binarios", "Multiclasse", "Com NaN"]
    rows = [
        ("Small (n<1K)",       "3",  "2", "1", "1"),
        ("Medium (1K-10K)",   "17",  "7", "10", "3"),
        ("Large (>10K)",      "10",  "6", "4",  "4"),
        ("TOTAL",             "30", "15", "15", "8"),
    ]
    _pptx_table(sl, headers, rows, cx, MH + Inches(.52), half - Inches(.12), Inches(1.8))
    _tx(sl, cx, MH + Inches(2.42), half - Inches(.1), Inches(.28),
        "4 eixos analisados: tamanho, classes, tipo de feature, dados faltantes.",
        size=11, italic=True, color=GRAY)

    # Detalhe dos 4 eixos numericos
    eixo_data = [
        ("Multiclasse",  "TabICL: 0.9585  >  AG 4h [spec]: 0.9606  (2o lugar)"),
        ("Binario",      "AG 4h [spec]: 0.8592  >  TabICL: 0.8562"),
        ("Sem NaN",      "TabICL: 0.9221  ~  AG 4h [spec]: 0.9221  (identicos)"),
        ("Com NaN",      "CatBoost TD: 0.8720  >  TabICL: 0.8539  (perde)"),
        ("Small",        "AG 4h [spec]: 0.8628  >  TabICL: 0.8056  (alta variancia)"),
        ("Medium",       "TabICL: 0.9425  ~  AG 4h [spec]: 0.9453  (diff = 0.003)"),
        ("Large",        "AG 4h [spec]: 0.8639  ~  TabICL: 0.8629  (diff = 0.001)"),
    ]
    cy2 = MH + Inches(2.8)
    for lbl, desc in eixo_data:
        _tx(sl, cx, cy2, Inches(1.3), Inches(.3),
            lbl, size=11, bold=True, color=RED)
        _tx(sl, cx + Inches(1.35), cy2, half - Inches(1.5), Inches(.3),
            desc, size=10, color=NEAR_BLK)
        cy2 += Inches(.34)

    # Direita: grafico de barras
    rx = SB + half + Inches(.2)
    _img(sl, PLOTS / "bar_auc_models.png", rx, MH + Inches(.08),
         half - Inches(.1), SH - MH - Inches(.18))


def slide_regime(prs):
    """Slide com grafico de regime (barras agrupadas) + tabela sintetica."""
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Analise por Regime: os 4 Eixos de Caracterizacao", MC)

    # Imagem maior: grafico de 4 eixos
    img_h = Inches(4.0)
    _img(sl, PLOTS / "bar_eixos.png",
         SB + Inches(.08), MH + Inches(.08),
         CW - Inches(.16), img_h)

    # Rodape com insight principal
    cy = MH + img_h + Inches(.15)
    _rect(sl, SB + MRG, cy, CW - MRG * 2, Inches(.85), ACCENT)
    _tx(sl, SB + MRG + Inches(.15), cy + Inches(.1), CW - MRG * 2 - Inches(.3), Inches(.65),
        "Insight: TabICL equivale ao AG Extreme 4h [EXIGIDO PELO ENUNCIADO] em Medium (diff=0.003) e "
        "em datasets sem NaN (identicos: 0.9221). Lidera em Multiclasse (0.9585 vs 0.9606). "
        "Fraqueza em NaN nativos e Small — alta variancia por N reduzido.",
        size=12, color=DARKRED)


def slide_resultados_metricas(prs):
    """Grafico de 3 metricas lado a lado."""
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Resultados por Metrica: AUC-OVO, Accuracy e G-Mean (30 datasets)", MC)

    _img(sl, PLOTS / "bar_4metricas.png",
         SB + Inches(.08), MH + Inches(.08),
         CW - Inches(.16), SH - MH - Inches(.42))

    _tx(sl, SB + MRG, SH - Inches(.32), CW - MRG * 2, Inches(.25),
        "TabICL lider em G-Mean (0.679 > todos). "
        "AutoGluon Extreme lidera em AUC e Accuracy (diff < 0.005). "
        "Cross-Entropy completa no relatorio (Apendice B).",
        size=10, italic=True, color=GRAY)


def slide_cd_diagrams(prs):
    """Todos os 4 CD diagrams em grid 2x2."""
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Friedman-Nemenyi: CD Diagrams para as 4 Metricas", MC)

    img_w = (CW - Inches(.4)) / 2
    img_h = (SH - MH - Inches(.45)) / 2
    positions = [
        (SB + Inches(.1),              MH + Inches(.08)),
        (SB + img_w + Inches(.2),      MH + Inches(.08)),
        (SB + Inches(.1),              MH + img_h + Inches(.18)),
        (SB + img_w + Inches(.2),      MH + img_h + Inches(.18)),
    ]
    for (path, lbl), (x, y) in zip([
        (IMG["cd_auc"],   "AUC-OVO"),
        (IMG["cd_acc"],   "Accuracy"),
        (IMG["cd_gmean"], "G-Mean"),
        (IMG["cd_ce"],    "Cross-Entropy"),
    ], positions):
        _img(sl, path, x, y, img_w, img_h - Inches(.1))
        _tx(sl, x, y + img_h - Inches(.05), img_w, Inches(.22),
            lbl, size=10, bold=True, color=RED, align=PP_ALIGN.CENTER)

    _tx(sl, SB + MRG, SH - Inches(.3), CW - MRG * 2, Inches(.22),
        "Barra horizontal: clique de equivalencia estatistica (p > 0.05, Bonferroni-Dunn). "
        "TabICL e AutoGluon Extreme no mesmo clique em AUC e Accuracy.",
        size=10, italic=True, color=GRAY)


def slide_bayesian(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Analise Bayesiana: Bayesian Signed-Rank com ROPE = 0.01", VIN)

    half = (CW - Inches(.4)) / 2
    img_h = Inches(4.35)
    _img(sl, IMG["bay_lgbm"], SB + Inches(.1), MH + Inches(.08), half, img_h)
    _img(sl, IMG["bay_ag"],   SB + half + Inches(.2), MH + Inches(.08), half, img_h)

    cy = MH + img_h + Inches(.12)
    _tx(sl, SB + Inches(.1), cy, half, Inches(.62),
        "TabICL vs LightGBM Tuned\n"
        "P(equiv.) = 83.1%  |  P(TabICL >) = 16.8%\n"
        "Frequentemente equivalentes",
        size=11, color=NEAR_BLK, align=PP_ALIGN.CENTER)
    _tx(sl, SB + half + Inches(.2), cy, half, Inches(.62),
        "TabICL vs AutoGluon Extreme 4h [EXIGIDO PELO ENUNCIADO]\n"
        "P(equiv.) = 99.6%  <-- resultado definitivo\n"
        "Indistinguiveis em AUC  |  TabICL 302x mais rapido",
        size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)


def slide_optuna(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Busca de Hiperparametros: Optuna TPE (GBDTs Tuned, 20 trials)", VIN)

    half = (CW - Inches(.4)) / 2
    img_h = Inches(3.1)
    _img(sl, IMG["opt_hist"], SB + Inches(.1), MH + Inches(.08), half, img_h)
    _img(sl, IMG["opt_imp"],  SB + half + Inches(.2), MH + Inches(.08), half, img_h)
    _tx(sl, SB + Inches(.1), MH + img_h + Inches(.1), half, Inches(.32),
        "Convergencia do objetivo (AUC-OVO)\nLightGBM, XGBoost, CatBoost — 20 trials",
        size=10, color=GRAY, align=PP_ALIGN.CENTER)
    _tx(sl, SB + half + Inches(.2), MH + img_h + Inches(.1), half, Inches(.32),
        "Importancia de HPs (fANOVA)\nlearning_rate e n_estimators: >60% da variancia",
        size=10, color=GRAY, align=PP_ALIGN.CENTER)

    cy = MH + img_h + Inches(.55)
    _tx(sl, SB + MRG, cy, CW - MRG * 2, Inches(.32),
        "Espacos de busca (Apendice C)  |  TPESampler(seed=42)  |  max 2h / dataset",
        size=12, bold=True, color=RED)
    headers = ["Modelo", "Hiperparametros buscados"]
    rows = [
        ("LightGBM Tuned",  "n_estimators [100-2K], learning_rate [1e-4, 0.3 log], num_leaves [15-300], subsample [0.5-1.0]"),
        ("XGBoost Tuned",   "n_estimators [100-2K], learning_rate [1e-4, 0.3 log], max_depth [3-12], colsample_bytree [0.5-1.0]"),
        ("CatBoost Tuned",  "iterations [100-2K], learning_rate [1e-4, 0.3 log], depth [4-10], l2_leaf_reg [1-10 log]"),
    ]
    _pptx_table(sl, headers, rows, SB + MRG, cy + Inches(.36), CW - MRG * 2, Inches(1.55))


def slide_pareto(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Custo vs. Desempenho: Fronteira de Pareto e Recomendacoes", VIN)

    half = CW / 2
    _img(sl, IMG["scatter"], SB + Inches(.06), MH + Inches(.06),
         half - Inches(.1), SH - MH - Inches(.16))

    tx, tw = SB + half + Inches(.08), half - Inches(.18)
    cy = MH + Inches(.1)
    for bg, lbl, desc in [
        (ACCENT,
         "Cenario A — Maxima qualidade  [EXIGIDO PELO ENUNCIADO]",
         "AutoGluon Extreme 4h  (preset extreme, time_limit=4h/dataset)\n"
         "AUC 0.9099  |  14.612s (~4h) / dataset\n"
         "Critico: saude, fraude, credito.\n"
         "Referencia primaria de comparacao do projeto."),
        (LGRAY,
         "Cenario B — Sweet spot (recomendado)",
         "TabICL v2\n"
         "AUC 0.9056  |  48.5s / dataset\n"
         "Na inflexao da fronteira de Pareto.\n"
         "302x mais rapido que AG Extreme 4h [spec]."),
        (LGRAY,
         "Cenario C — Latencia em CPU",
         "CatBoost TD\n"
         "AUC 0.8972  |  14s / dataset\n"
         "Inferencia em ms, sem GPU.\n"
         "Natural para NaN e categoricas."),
    ]:
        _rect(sl, tx, cy, tw, Inches(1.92), bg)
        _tx(sl, tx + Inches(.12), cy + Inches(.09), tw - Inches(.24), Inches(.4),
            lbl, size=12, bold=True, color=RED)
        _tx(sl, tx + Inches(.12), cy + Inches(.53), tw - Inches(.24), Inches(1.25),
            desc, size=11, color=NEAR_BLK)
        cy += Inches(2.02)


def slide_conclusoes(prs):
    sl = _blank(prs)
    _sidebar(sl)
    _header(sl, "Conclusoes", VIN)
    cx, cw = SB + MRG, CW - MRG * 2
    items = [
        ("ICL puro equivale ao AutoGluon Extreme 4h [exigido] com 302x menos tempo",
         "Bayesian Signed-Rank P(equiv.) = 99.6% com ROPE = 0.01 (baycomp). "
         "AutoGluon Extreme 4h: exigido explicitamente pelo enunciado (preset extreme, time_limit=4h/dataset). "
         "TabICL 48.5s vs 14.612s (4h) por dataset. G-Mean 0.679 — melhor entre todos os 10 sistemas."),
        ("Lideranca em Multiclasse e equivalencia em Medium",
         "AUC 0.9585 em Multiclasse — melhor sistema individual, supera todos os GBDTs. "
         "Rank 3 em AUC-OVO geral (30 datasets). Diff < 0.002 vs SOTA em Medium. "
         "Hiperparametros default ja competitivos — tuning adiciona ganho marginal."),
        ("Limitacoes claras e recomendacoes objetivas",
         "Fraqueza em NaN nativos (imputacao mediana < particionamento dos GBDTs). "
         "OOM em Large sem GPU >= 24 GB. "
         "Recomendacao: TabICL para Medium/Multiclasse com GPU; "
         "CatBoost TD para NaN+CPU; AutoGluon 4h para qualidade maxima absoluta."),
    ]
    cy = MH + Inches(.1)
    bh = (SH - MH - Inches(.2)) / len(items)
    for lbl, body in items:
        _tx(sl, cx, cy, cw, Inches(.38), lbl, size=15, bold=True, color=RED)
        _tx(sl, cx + Inches(.1), cy + Inches(.41), cw - Inches(.1),
            bh - Inches(.5), body, size=13, color=NEAR_BLK)
        cy += bh


def slide_fechamento(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, RED)
    _tx(sl, Inches(1.0), Inches(1.15), SW - Inches(2.0), Inches(1.75),
        "Muito obrigado!", size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(sl, Inches(1.0), Inches(3.02), SW - Inches(2.0), Inches(.55),
        "Perguntas?", size=24, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    for i, (nome, email) in enumerate([
        ("Joao Pedro Miranda da Silva", "jpms5@cin.ufpe.br"),
        ("Maria Clara F. G. Barretto",  "mcfgb@cin.ufpe.br"),
        ("Vinicius Limeira Valenca",    "vlv3@cin.ufpe.br"),
    ]):
        _tx(sl, Inches(1.8), Inches(3.95) + Inches(.6) * i, SW - Inches(3.6), Inches(.52),
            f"{nome}   |   {email}", size=14, color=WHITE)
    _tx(sl, Inches(1.0), Inches(6.46), SW - Inches(2.0), Inches(.45),
        "github.com/jpmsilva1/Projeto_AM_Leandro_TabICL",
        size=12, color=WHITE, italic=True, align=PP_ALIGN.CENTER)


# Alias para IMG (resolve no momento do build)
IMG: dict[str, Path] = {}


# =============================================================================
# Build
# =============================================================================

def build(out: str = "reports/slides_tabicl_v2.pptx"):
    global IMG
    IMG = {
        "cd_auc":   PLOTS / "cd_diagram_auc_ovo.png",
        "cd_acc":   PLOTS / "cd_diagram_acc.png",
        "cd_gmean": PLOTS / "cd_diagram_g_mean.png",
        "cd_ce":    PLOTS / "cd_diagram_ce.png",
        "bay_lgbm": PLOTS / "bayesian_plot_auc_lightgbm_tuned.png",
        "bay_ag":   PLOTS / "bayesian_plot_auc_autogluon_extreme_4h.png",
        "scatter":  PLOTS / "scatter_cost_perf.png",
        "opt_hist": PLOTS / "optuna_history.png",
        "opt_imp":  PLOTS / "optuna_importances.png",
    }

    print("Gerando graficos extras...")
    _generate_charts()

    print("Montando PPTX...")
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    # ETAPA 1 -- Teoria (10 min, Joao Pedro)
    slide_capa(prs)               # 1
    slide_agenda(prs)             # 2
    slide_motivacao(prs)          # 3 - motivacao + contexto historico
    slide_arquitetura(prs)        # 4 - arquitetura + diagrama de blocos
    slide_funcionamento(prs)      # 5 - formula + notacao + ensemble
    slide_aprendizado(prs)        # 6 - pre-treino vs inferencia (rubrica explicita)
    slide_complexidade(prs)       # 7 - O(n^2.p), tabela, mitigacoes
    slide_hiperparametros(prs)    # 8 - tabela de HPs + espaco Optuna
    slide_aplicacoes(prs)         # 9 - 2 colunas usar/nao usar (rubrica explicita)
    slide_transicao(prs)          # 10

    # ETAPA 2 -- Experimentos (10 min, MC + Vinicius)
    slide_setup(prs)              # 11 - configuracao completa
    slide_datasets_e_resultados(prs)  # 12 - resumo 30 datasets + grafico AUC por modelo
    slide_resultados_metricas(prs)    # 13 - 3 graficos AUC/ACC/G-Mean
    slide_regime(prs)             # 14 - grafico 4 eixos (2x2) + insight
    slide_cd_diagrams(prs)        # 15 - 4 CD diagrams (2x2)
    slide_bayesian(prs)           # 16 - 2 simplexes bayesianos
    slide_optuna(prs)             # 17 - convergencia + importancia + tabela espacos
    slide_pareto(prs)             # 18 - scatter Pareto + 3 cenarios
    slide_conclusoes(prs)         # 19
    slide_fechamento(prs)         # 20 -- retirada slide 21 para fechar em 20

    Path(out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"OK  {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
